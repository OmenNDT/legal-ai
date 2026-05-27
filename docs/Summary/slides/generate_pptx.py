#!/usr/bin/env python3
"""Generate PowerPoint presentation for Hybrid Summarization on CUAD."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# ── Colors ──
BG_DARK    = RGBColor(0x0A, 0x0E, 0x1A)
BG_CARD    = RGBColor(0x1A, 0x22, 0x36)
ACCENT     = RGBColor(0x3B, 0x82, 0xF6)
ACCENT_L   = RGBColor(0x60, 0xA5, 0xFA)
GREEN      = RGBColor(0x10, 0xB9, 0x81)
ORANGE     = RGBColor(0xF5, 0x9E, 0x0B)
RED        = RGBColor(0xEF, 0x44, 0x44)
PURPLE     = RGBColor(0x8B, 0x5C, 0xF6)
CYAN       = RGBColor(0x06, 0xB6, 0xD4)
WHITE      = RGBColor(0xF1, 0xF5, 0xF9)
GRAY       = RGBColor(0x94, 0xA3, 0xB8)
MUTED      = RGBColor(0x64, 0x74, 0x8B)
BORDER     = RGBColor(0x1E, 0x29, 0x3B)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H

# Use blank layout
blank_layout = prs.slide_layouts[6]


def set_bg(slide, color=BG_DARK):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_textbox(slide, left, top, width, height, text, font_size=14,
                color=WHITE, bold=False, alignment=PP_ALIGN.LEFT, font_name='Calibri'):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_para(text_frame, text, font_size=14, color=GRAY, bold=False,
             alignment=PP_ALIGN.LEFT, space_before=Pt(4), font_name='Calibri'):
    p = text_frame.add_paragraph()
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    if space_before:
        p.space_before = space_before
    return p


def add_card(slide, left, top, width, height, border_color=ACCENT):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = BG_CARD
    shape.line.color.rgb = border_color
    shape.line.width = Pt(1.5)
    shape.shadow.inherit = False
    return shape


def add_slide_number(slide, num, total=12):
    add_textbox(slide, Inches(11.8), Inches(0.3), Inches(1.2), Inches(0.35),
                f"{num:02d} / {total}", font_size=11, color=MUTED,
                alignment=PP_ALIGN.RIGHT, font_name='Consolas')


def add_tag(slide, left, top, text, color=ACCENT_L):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, Inches(1.8), Inches(0.32))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0x1E, 0x2A, 0x42)
    shape.line.color.rgb = color
    shape.line.width = Pt(1)
    tf = shape.text_frame
    tf.paragraphs[0].text = text
    tf.paragraphs[0].font.size = Pt(9)
    tf.paragraphs[0].font.color.rgb = color
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.name = 'Calibri'
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE


def add_stat_card(slide, left, top, value, unit, label, val_color=ACCENT_L):
    card = add_card(slide, left, top, Inches(2.6), Inches(1.3))
    tf = card.text_frame
    tf.word_wrap = True
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    tf.paragraphs[0].text = value
    tf.paragraphs[0].font.size = Pt(36)
    tf.paragraphs[0].font.color.rgb = val_color
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.name = 'Calibri'
    p2 = tf.add_paragraph()
    p2.text = unit
    p2.font.size = Pt(12)
    p2.font.color.rgb = GRAY
    p2.alignment = PP_ALIGN.CENTER
    p3 = tf.add_paragraph()
    p3.text = label
    p3.font.size = Pt(11)
    p3.font.color.rgb = MUTED
    p3.alignment = PP_ALIGN.CENTER
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE


# ════════════════════════════════════════════════════════════════
# SLIDE 1 — Title + Problem Intro
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
set_bg(slide)
add_slide_number(slide, 1)
add_tag(slide, Inches(4.8), Inches(1.0), "BÀI TOÁN & GIỚI THIỆU", ACCENT_L)

add_textbox(slide, Inches(1.5), Inches(1.5), Inches(10.3), Inches(1.2),
            "Hệ thống Tóm tắt Văn bản Lai", font_size=40, color=WHITE, bold=True,
            alignment=PP_ALIGN.CENTER)
add_textbox(slide, Inches(1.5), Inches(2.4), Inches(10.3), Inches(0.6),
            "Hybrid Summarization trên CUAD", font_size=32, color=ACCENT_L, bold=True,
            alignment=PP_ALIGN.CENTER)
add_textbox(slide, Inches(2.5), Inches(3.1), Inches(8.3), Inches(0.5),
            "Kết hợp Extractive + Abstractive để tóm tắt hợp đồng pháp lý dài",
            font_size=16, color=GRAY, alignment=PP_ALIGN.CENTER)

# Stats row
add_stat_card(slide, Inches(1.2), Inches(3.9), "510", "", "Hợp đồng pháp lý", ACCENT_L)
add_stat_card(slide, Inches(4.1), Inches(3.9), "7.8K", "từ", "Trung bình / file", ORANGE)
add_stat_card(slide, Inches(7.0), Inches(3.9), "47.7K", "từ", "Tối đa (~60 trang)", RED)
add_stat_card(slide, Inches(9.9), Inches(3.9), "83", "cột", "Clause tham chiếu", PURPLE)

# Input / Output cards
card_in = add_card(slide, Inches(2.0), Inches(5.5), Inches(4.2), Inches(1.4), ACCENT)
tf = card_in.text_frame
tf.word_wrap = True
tf.paragraphs[0].text = "📥  Input"
tf.paragraphs[0].font.size = Pt(16)
tf.paragraphs[0].font.color.rgb = ACCENT_L
tf.paragraphs[0].font.bold = True
p = tf.add_paragraph()
p.text = "Hợp đồng pháp lý tiếng Anh thô (raw text) dài 7K–47K từ, không có bản tóm tắt sẵn."
p.font.size = Pt(12)
p.font.color.rgb = GRAY

card_out = add_card(slide, Inches(7.1), Inches(5.5), Inches(4.2), Inches(1.4), GREEN)
tf = card_out.text_frame
tf.word_wrap = True
tf.paragraphs[0].text = "📤  Output"
tf.paragraphs[0].font.size = Pt(16)
tf.paragraphs[0].font.color.rgb = GREEN
tf.paragraphs[0].font.bold = True
p = tf.add_paragraph()
p.text = "Bản tóm tắt ngắn 150–300 từ, mượt mà, giữ nguyên thông tin cốt lõi của hợp đồng."
p.font.size = Pt(12)
p.font.color.rgb = GRAY


# ════════════════════════════════════════════════════════════════
# SLIDE 2 — Thách thức & Khó khăn
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
set_bg(slide)
add_slide_number(slide, 2)
add_tag(slide, Inches(0.8), Inches(0.5), "THÁCH THỨC & KHÓ KHĂN", RED)

add_textbox(slide, Inches(0.8), Inches(1.0), Inches(11), Inches(0.6),
            "Tại sao không thể tóm tắt trực tiếp?", font_size=28, color=WHITE, bold=True)

challenges = [
    ("⚠️ Giới hạn token mô hình", RED,
     "BART/T5 chỉ nhận ≤ 1024 token đầu vào. Hợp đồng dài nhất ~47K từ ≈ 60K token → vượt xa 60 lần giới hạn."),
    ("⚠️ Cắt cụt = Mất ngữ cảnh", ORANGE,
     "Nếu chỉ lấy 1024 token đầu → bỏ lỡ điều khoản quan trọng ở cuối hợp đồng (phụ lục, bảo mật, chấm dứt)."),
    ("⚠️ Tốn tài nguyên GPU", PURPLE,
     "Mô hình Deep Learning xử lý văn bản dài rất tốn RAM/VRAM. Không khả thi với tài nguyên hạn chế."),
    ("⚠️ Không có Gold Summary", CYAN,
     "CUAD không cung cấp bản tóm tắt mẫu. Chỉ có master_clauses.csv (83 cột clause) do luật sư gán nhãn — phải tận dụng làm reference."),
]

for i, (title, color, desc) in enumerate(challenges):
    col = i % 2
    row = i // 2
    left = Inches(0.8 + col * 6.1)
    top = Inches(1.8 + row * 2.6)
    card = add_card(slide, left, top, Inches(5.7), Inches(2.2), color)
    tf = card.text_frame
    tf.word_wrap = True
    tf.paragraphs[0].text = title
    tf.paragraphs[0].font.size = Pt(16)
    tf.paragraphs[0].font.color.rgb = color
    tf.paragraphs[0].font.bold = True
    p = tf.add_paragraph()
    p.text = desc
    p.font.size = Pt(12)
    p.font.color.rgb = GRAY
    p.space_before = Pt(8)


# ════════════════════════════════════════════════════════════════
# SLIDE 3 — Tổng quát Data Flow
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
set_bg(slide)
add_slide_number(slide, 3)
add_tag(slide, Inches(0.8), Inches(0.5), "DATA FLOW", ACCENT_L)

add_textbox(slide, Inches(0.8), Inches(1.0), Inches(11), Inches(0.6),
            "Luồng dữ liệu tổng quát — Hybrid Pipeline", font_size=28, color=WHITE, bold=True)
add_textbox(slide, Inches(0.8), Inches(1.6), Inches(10), Inches(0.4),
            "4 bước tuần tự: Preprocess → Extractive → Abstractive → Evaluate",
            font_size=14, color=GRAY)

# Flow steps
steps = [
    ("①", "Preprocess", "Clean + Tách câu", ACCENT_L),
    ("②", "Extractive", "Lọc 20% câu cốt lõi", GREEN),
    ("③", "Abstractive", "BART viết lại", PURPLE),
    ("④", "Evaluate", "ROUGE + BERTScore", ORANGE),
]

for i, (num, label, desc, color) in enumerate(steps):
    left = Inches(0.8 + i * 3.1)
    # Circle
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, left + Inches(0.7), Inches(2.3), Inches(0.9), Inches(0.9))
    shape.fill.solid()
    shape.fill.fore_color.rgb = BG_CARD
    shape.line.color.rgb = color
    shape.line.width = Pt(2)
    tf = shape.text_frame
    tf.paragraphs[0].text = num
    tf.paragraphs[0].font.size = Pt(22)
    tf.paragraphs[0].font.color.rgb = color
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE

    add_textbox(slide, left, Inches(3.4), Inches(2.5), Inches(0.4),
                label, font_size=16, color=color, bold=True, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, left, Inches(3.8), Inches(2.5), Inches(0.3),
                desc, font_size=12, color=MUTED, alignment=PP_ALIGN.CENTER)

    # Arrow
    if i < 3:
        add_textbox(slide, left + Inches(2.5), Inches(2.55), Inches(0.6), Inches(0.4),
                    "→", font_size=24, color=MUTED, alignment=PP_ALIGN.CENTER)

# Detail cards
details = [
    ("① PREPROCESS", "ContractLoader → TextCleaner → SentenceSplitter",
     "Chuẩn hoá Unicode NFKC, loại noise (Page X of Y, Source:), tách câu NLTK punkt, lọc câu <5 hoặc >80 từ", ACCENT),
    ("② EXTRACTIVE", "TF-IDF / TextRank / KMeans / Ensemble",
     "Chấm điểm câu → lấy top 20% → giữ thứ tự gốc. Đảm bảo ≤ 1024 token cho BART", GREEN),
    ("③ ABSTRACTIVE", "LongDocChunker → BartSummarizer",
     "Chunk ≤ 1024 token (overlap 50), beam=4, hierarchical nếu vẫn dài. Output: 150–300 từ", PURPLE),
    ("④ EVALUATE", "ReferenceBuilder → ROUGE + BERTScore",
     "Gộp clauses từ CSV thành \"tóm tắt vàng\". ROUGE-1/2/L + BERTScore (roberta-large)", ORANGE),
]

for i, (step, title, desc, color) in enumerate(details):
    col = i % 2
    row = i // 2
    left = Inches(0.8 + col * 6.1)
    top = Inches(4.4 + row * 1.5)
    card = add_card(slide, left, top, Inches(5.7), Inches(1.3), color)
    tf = card.text_frame
    tf.word_wrap = True
    tf.paragraphs[0].text = step
    tf.paragraphs[0].font.size = Pt(11)
    tf.paragraphs[0].font.color.rgb = color
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.name = 'Consolas'
    p = tf.add_paragraph()
    p.text = title
    p.font.size = Pt(13)
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.space_before = Pt(4)
    p = tf.add_paragraph()
    p.text = desc
    p.font.size = Pt(10)
    p.font.color.rgb = GRAY
    p.space_before = Pt(4)


# ════════════════════════════════════════════════════════════════
# SLIDE 4 — Preprocess chi tiết
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
set_bg(slide)
add_slide_number(slide, 4)
add_tag(slide, Inches(0.8), Inches(0.5), "BƯỚC ①", ACCENT_L)

add_textbox(slide, Inches(0.8), Inches(1.0), Inches(11), Inches(0.6),
            "Preprocess — Từ raw text → danh sách câu sạch", font_size=28, color=WHITE, bold=True)

# Left column - 3 cards
preprocess_items = [
    ("📄 ContractLoader", "Đọc file .txt → object Contract{doc_id, raw_text, word_count}", ACCENT_L),
    ("🧹 TextCleaner", "• Chuẩn hoá Unicode NFKC\n• Loại Page X of Y, dòng Source:\n• Gộp khoảng trắng / xuống dòng dư thừa", ACCENT_L),
    ("✂️ SentenceSplitter", "• Ưu tiên nltk.sent_tokenize (punkt)\n• Fallback regex nếu NLTK lỗi\n• Lọc câu <5 từ (tiêu đề) hoặc >80 từ (nhiễu)", ACCENT_L),
]

for i, (title, desc, color) in enumerate(preprocess_items):
    top = Inches(1.8 + i * 1.7)
    card = add_card(slide, Inches(0.8), top, Inches(5.5), Inches(1.5), color)
    tf = card.text_frame
    tf.word_wrap = True
    tf.paragraphs[0].text = title
    tf.paragraphs[0].font.size = Pt(15)
    tf.paragraphs[0].font.color.rgb = color
    tf.paragraphs[0].font.bold = True
    p = tf.add_paragraph()
    p.text = desc
    p.font.size = Pt(11)
    p.font.color.rgb = GRAY
    p.space_before = Pt(6)

# Right column - flow code
card = add_card(slide, Inches(6.8), Inches(1.8), Inches(5.7), Inches(5.1), PURPLE)
tf = card.text_frame
tf.word_wrap = True
flow_text = """Raw: "AGREEMENT ... Page 1 of 15 ..."
  │
  ▼ ContractLoader
Contract { doc_id, raw_text, word_count: 7861 }
  │
  ▼ TextCleaner
"AGREEMENT ..."  ← (đã bỏ Page/Source)
  │
  ▼ SentenceSplitter
List[Sentence] = [
  (0, "AGREEMENT made as of...", 6),
  (1, "The Company shall...", 12),
  ...
]
  │
  ▼ Lọc 5–80 từ
List[Sentence] — chỉ giữ câu có ý nghĩa"""
tf.paragraphs[0].text = "🔄 Luồng xử lý"
tf.paragraphs[0].font.size = Pt(14)
tf.paragraphs[0].font.color.rgb = WHITE
tf.paragraphs[0].font.bold = True
p = tf.add_paragraph()
p.text = flow_text
p.font.size = Pt(10)
p.font.color.rgb = GRAY
p.font.name = 'Consolas'
p.space_before = Pt(8)


# ════════════════════════════════════════════════════════════════
# SLIDE 5 — Extractive: 3 thuật toán
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
set_bg(slide)
add_slide_number(slide, 5)
add_tag(slide, Inches(0.8), Inches(0.5), "BƯỚC ②", GREEN)

add_textbox(slide, Inches(0.8), Inches(1.0), Inches(11), Inches(0.6),
            "Extractive — 3 thuật toán + Ensemble", font_size=28, color=WHITE, bold=True)
add_textbox(slide, Inches(0.8), Inches(1.6), Inches(10), Inches(0.4),
            "Chấm điểm câu → lấy top 20% → giữ thứ tự gốc để bản tóm tắt mạch lạc",
            font_size=14, color=GRAY)

algos = [
    ("TF-IDF", ACCENT_L,
     "Mỗi câu = 1 \"document\". TfidfVectorizer(ngram=(1,2)) → ma trận TF-IDF.",
     "score = Σ TF-IDF(câu) / |từ|",
     "Câu chứa nhiều từ khoá quan trọng & hiếm → điểm cao"),
    ("TextRank ★ ưu tiên", GREEN,
     "Cosine similarity giữa các câu → đồ thị có trọng số → pagerank(damping=0.85)",
     "PageRank — câu nào được \"nhiều câu vote\" thì quan trọng",
     "Thường mạnh nhất → trọng số 1.5 trong ensemble"),
    ("K-Means", PURPLE,
     "Embed câu bằng all-MiniLM-L6-v2 (384 chiều) → gom K cụm.",
     "score = 1 / (1 + dist_to_centroid)",
     "Mỗi cụm = 1 chủ đề → lấy đại diện → tóm tắt đa chủ đề"),
]

for i, (name, color, desc, formula, intuition) in enumerate(algos):
    left = Inches(0.8 + i * 4.1)
    card = add_card(slide, left, Inches(2.2), Inches(3.7), Inches(3.8), color)
    tf = card.text_frame
    tf.word_wrap = True
    tf.paragraphs[0].text = name
    tf.paragraphs[0].font.size = Pt(16)
    tf.paragraphs[0].font.color.rgb = color
    tf.paragraphs[0].font.bold = True
    p = tf.add_paragraph()
    p.text = desc
    p.font.size = Pt(11)
    p.font.color.rgb = GRAY
    p.space_before = Pt(8)
    p = tf.add_paragraph()
    p.text = formula
    p.font.size = Pt(11)
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.font.name = 'Consolas'
    p.space_before = Pt(10)
    p = tf.add_paragraph()
    p.text = intuition
    p.font.size = Pt(10)
    p.font.color.rgb = MUTED
    p.font.italic = True
    p.space_before = Pt(8)

# Ensemble bar
card = add_card(slide, Inches(0.8), Inches(6.2), Inches(11.7), Inches(0.9), ORANGE)
tf = card.text_frame
tf.word_wrap = True
tf.paragraphs[0].text = "Ensemble:  Min-max normalize → cộng theo trọng số    TF-IDF × 1.0  +  TextRank × 1.5  +  KMeans × 1.0"
tf.paragraphs[0].font.size = Pt(13)
tf.paragraphs[0].font.color.rgb = ORANGE
tf.paragraphs[0].font.bold = True
tf.paragraphs[0].font.name = 'Consolas'
tf.vertical_anchor = MSO_ANCHOR.MIDDLE


# ════════════════════════════════════════════════════════════════
# SLIDE 6 — Abstractive (BART)
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
set_bg(slide)
add_slide_number(slide, 6)
add_tag(slide, Inches(0.8), Inches(0.5), "BƯỚC ③", PURPLE)

add_textbox(slide, Inches(0.8), Inches(1.0), Inches(11), Inches(0.6),
            "Abstractive — BART-large-CNN viết lại", font_size=28, color=WHITE, bold=True)

# Left column
items = [
    ("📦 LongDocChunker", "Nếu câu sau extractive vẫn > 1024 token → chia chunk ≤ 1024 token, overlap 50 token giữ ngữ cảnh.", PURPLE),
    ("🤖 BartSummarizer", "• Model: facebook/bart-large-cnn\n• Beam search: num_beams=4\n• length_penalty=2.0, no_repeat_ngram_size=3\n• Max output: 256 token → ~150–300 từ", PURPLE),
    ("🔄 Hierarchical", "Nhiều chunk → ghép tóm tắt → nếu tổng vẫn > 1024 → tóm tắt thêm 1 vòng nữa.", PURPLE),
]

for i, (title, desc, color) in enumerate(items):
    top = Inches(1.8 + i * 1.7)
    card = add_card(slide, Inches(0.8), top, Inches(5.5), Inches(1.5), color)
    tf = card.text_frame
    tf.word_wrap = True
    tf.paragraphs[0].text = title
    tf.paragraphs[0].font.size = Pt(15)
    tf.paragraphs[0].font.color.rgb = color
    tf.paragraphs[0].font.bold = True
    p = tf.add_paragraph()
    p.text = desc
    p.font.size = Pt(11)
    p.font.color.rgb = GRAY
    p.space_before = Pt(6)

# Right column - chunking flow
card = add_card(slide, Inches(6.8), Inches(1.8), Inches(5.7), Inches(5.1), PURPLE)
tf = card.text_frame
tf.word_wrap = True
chunk_flow = """[Câu extractive — 47K từ → ~9400 token]
  │
  ▼ LongDocChunker
  ├── Chunk 1 (≤1024 token, overlap 50)
  ├── Chunk 2 (≤1024 token, overlap 50)
  ├── Chunk 3 ...
  └── Chunk N
  │
  ▼ BartSummarizer (mỗi chunk)
  ├── Summary 1 (~150 từ)
  ├── Summary 2 (~150 từ)
  └── ...
  │
  ▼ Ghép summaries
  Tổng > 1024 token?
    ├── YES → tóm tắt thêm 1 vòng (hierarchical)
    └── NO  → bản tóm tắt cuối cùng ✅"""
tf.paragraphs[0].text = "Chunking & Hierarchical"
tf.paragraphs[0].font.size = Pt(14)
tf.paragraphs[0].font.color.rgb = WHITE
tf.paragraphs[0].font.bold = True
p = tf.add_paragraph()
p.text = chunk_flow
p.font.size = Pt(10)
p.font.color.rgb = GRAY
p.font.name = 'Consolas'
p.space_before = Pt(8)


# ════════════════════════════════════════════════════════════════
# SLIDE 7 — Evaluate
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
set_bg(slide)
add_slide_number(slide, 7)
add_tag(slide, Inches(0.8), Inches(0.5), "BƯỚC ④", ORANGE)

add_textbox(slide, Inches(0.8), Inches(1.0), Inches(11), Inches(0.6),
            "Evaluate — \"Tóm tắt vàng\" từ CUAD + Đánh giá", font_size=28, color=WHITE, bold=True)

# Left
card = add_card(slide, Inches(0.8), Inches(1.8), Inches(5.5), Inches(2.2), ORANGE)
tf = card.text_frame
tf.word_wrap = True
tf.paragraphs[0].text = "🏆 ReferenceBuilder"
tf.paragraphs[0].font.size = Pt(16)
tf.paragraphs[0].font.color.rgb = ORANGE
tf.paragraphs[0].font.bold = True
p = tf.add_paragraph()
p.text = "Đọc master_clauses.csv (83 cột) → gộp các clause của từng doc (loại trùng) → \"tóm tắt vàng\" để chấm điểm."
p.font.size = Pt(12)
p.font.color.rgb = GRAY
p.space_before = Pt(8)
p = tf.add_paragraph()
p.text = "Điểm độc đáo: CUAD không có gold summary → tận dụng clause labels do luật sư gán nhãn làm reference"
p.font.size = Pt(11)
p.font.color.rgb = ORANGE
p.font.bold = True
p.space_before = Pt(10)

card = add_card(slide, Inches(0.8), Inches(4.2), Inches(5.5), Inches(2.8), GREEN)
tf = card.text_frame
tf.word_wrap = True
tf.paragraphs[0].text = "📊 RougeEvaluator"
tf.paragraphs[0].font.size = Pt(16)
tf.paragraphs[0].font.color.rgb = GREEN
tf.paragraphs[0].font.bold = True
for line in ["ROUGE-1: unigram overlap", "ROUGE-2: bigram overlap (fluency)", "ROUGE-L: longest common subsequence", "Mỗi metric: precision / recall / F1"]:
    p = tf.add_paragraph()
    p.text = "• " + line
    p.font.size = Pt(12)
    p.font.color.rgb = GRAY
    p.space_before = Pt(4)

# Right
card = add_card(slide, Inches(6.8), Inches(1.8), Inches(5.7), Inches(2.2), PURPLE)
tf = card.text_frame
tf.word_wrap = True
tf.paragraphs[0].text = "🧠 BertScoreEvaluator"
tf.paragraphs[0].font.size = Pt(16)
tf.paragraphs[0].font.color.rgb = PURPLE
tf.paragraphs[0].font.bold = True
p = tf.add_paragraph()
p.text = "Dùng roberta-large để đánh giá nghĩa (semantic similarity), không chỉ so trùng từ vựng."
p.font.size = Pt(12)
p.font.color.rgb = GRAY
p.space_before = Pt(8)

card = add_card(slide, Inches(6.8), Inches(4.2), Inches(5.7), Inches(2.8), ACCENT)
tf = card.text_frame
tf.word_wrap = True
tf.paragraphs[0].text = "📋 EvalRunner"
tf.paragraphs[0].font.size = Pt(16)
tf.paragraphs[0].font.color.rgb = ACCENT_L
tf.paragraphs[0].font.bold = True
p = tf.add_paragraph()
p.text = "Benchmark trên toàn bộ 510 file → lưu kết quả vào outputs/eval/rouge_*.json"
p.font.size = Pt(12)
p.font.color.rgb = GRAY
p.space_before = Pt(8)
p = tf.add_paragraph()
p.text = "POST /api/eval/run\n{ \"limit\": 510, \"extractor\": \"ensemble\" }"
p.font.size = Pt(10)
p.font.color.rgb = MUTED
p.font.name = 'Consolas'
p.space_before = Pt(10)


# ════════════════════════════════════════════════════════════════
# SLIDE 8 — Chiến lược Huấn luyện
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
set_bg(slide)
add_slide_number(slide, 8)
add_tag(slide, Inches(0.8), Inches(0.5), "CHIẾN LƯỢC HUẤN LUYỆN", PURPLE)

add_textbox(slide, Inches(0.8), Inches(1.0), Inches(11), Inches(0.6),
            "Fine-tune BART trên CUAD — RTX 3090 24GB", font_size=28, color=WHITE, bold=True)

# Left - Dataset
card = add_card(slide, Inches(0.8), Inches(1.8), Inches(5.5), Inches(2.0), PURPLE)
tf = card.text_frame
tf.word_wrap = True
tf.paragraphs[0].text = "📦 Dataset: CuadDatasetBuilder"
tf.paragraphs[0].font.size = Pt(15)
tf.paragraphs[0].font.color.rgb = PURPLE
tf.paragraphs[0].font.bold = True
for line in ["• Input: câu trích xuất bằng TextRank", "• Target: clauses gộp từ CSV (tóm tắt vàng)", "• Chia 80/10/10 train/val/test (seed=42)", "• Cache vào dataset_cuad.json → không build lại"]:
    p = tf.add_paragraph()
    p.text = line
    p.font.size = Pt(11)
    p.font.color.rgb = GRAY
    p.space_before = Pt(4)

# Left - Training params
card = add_card(slide, Inches(0.8), Inches(4.0), Inches(5.5), Inches(3.0), ACCENT)
tf = card.text_frame
tf.word_wrap = True
tf.paragraphs[0].text = "🏋️ BartFineTuner (Seq2SeqTrainer)"
tf.paragraphs[0].font.size = Pt(15)
tf.paragraphs[0].font.color.rgb = ACCENT_L
tf.paragraphs[0].font.bold = True

params = [("Epochs", "3"), ("Batch size", "2"), ("Grad accum", "8 → eff batch = 16"),
          ("Learning rate", "3e-5"), ("Precision", "fp16"), ("Beam width", "4")]
for k, v in params:
    p = tf.add_paragraph()
    p.text = f"  {k}: {v}"
    p.font.size = Pt(11)
    p.font.color.rgb = GRAY
    p.space_before = Pt(3)

# Right - LoRA
card = add_card(slide, Inches(6.8), Inches(1.8), Inches(5.7), Inches(1.8), ORANGE)
tf = card.text_frame
tf.word_wrap = True
tf.paragraphs[0].text = "⚡ LoRA (tuỳ chọn)"
tf.paragraphs[0].font.size = Pt(15)
tf.paragraphs[0].font.color.rgb = ORANGE
tf.paragraphs[0].font.bold = True
p = tf.add_paragraph()
p.text = "Cờ --use_lora → giảm VRAM bằng LoRA (r=16). Phù hợp khi GPU < 24GB."
p.font.size = Pt(12)
p.font.color.rgb = GRAY
p.space_before = Pt(8)

# Right - Deploy
card = add_card(slide, Inches(6.8), Inches(3.8), Inches(5.7), Inches(3.2), GREEN)
tf = card.text_frame
tf.word_wrap = True
tf.paragraphs[0].text = "🚀 Deploy qua Fabric"
tf.paragraphs[0].font.size = Pt(15)
tf.paragraphs[0].font.color.rgb = GREEN
tf.paragraphs[0].font.bold = True
deploy_text = """fab sync       # rsync code lên worker1
fab setup      # venv + torch CUDA 12.1
fab gpu        # check nvidia-smi
fab train --epochs=3 --batch=2 \\
  --grad-accum=8
fab tail       # xem log train
fab pull-model # kéo model về local"""
p = tf.add_paragraph()
p.text = deploy_text
p.font.size = Pt(10)
p.font.color.rgb = GRAY
p.font.name = 'Consolas'
p.space_before = Pt(8)


# ════════════════════════════════════════════════════════════════
# SLIDE 9 — Backend + Frontend
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
set_bg(slide)
add_slide_number(slide, 9)
add_tag(slide, Inches(0.8), Inches(0.5), "KIẾN TRÚC", ACCENT_L)

add_textbox(slide, Inches(0.8), Inches(1.0), Inches(11), Inches(0.6),
            "Backend Flask API (port 9020) + Frontend React", font_size=28, color=WHITE, bold=True)

# Left - API
card = add_card(slide, Inches(0.8), Inches(1.8), Inches(5.5), Inches(5.0), ACCENT)
tf = card.text_frame
tf.word_wrap = True
tf.paragraphs[0].text = "🔌 REST API Endpoints"
tf.paragraphs[0].font.size = Pt(15)
tf.paragraphs[0].font.color.rgb = ACCENT_L
tf.paragraphs[0].font.bold = True

endpoints = [
    ("GET", "/api/health", "Trạng thái + device"),
    ("GET", "/api/documents", "Danh sách doc (search, page)"),
    ("GET", "/api/documents/:id", "Nội dung 1 doc + ref"),
    ("POST", "/api/extract", "Chỉ extractive (nhanh)"),
    ("POST", "/api/summarize", "Pipeline lai đầy đủ"),
    ("POST", "/api/eval/run", "Benchmark toàn bộ"),
]
for method, path, desc in endpoints:
    p = tf.add_paragraph()
    run = p.add_run()
    run.text = f"  {method} "
    run.font.size = Pt(10)
    run.font.color.rgb = GREEN if method == "GET" else ORANGE
    run.font.bold = True
    run.font.name = 'Consolas'
    run2 = p.add_run()
    run2.text = f"{path}  — {desc}"
    run2.font.size = Pt(10)
    run2.font.color.rgb = GRAY
    run2.font.name = 'Consolas'
    p.space_before = Pt(6)

# Right - Frontend
card = add_card(slide, Inches(6.8), Inches(1.8), Inches(5.7), Inches(2.4), GREEN)
tf = card.text_frame
tf.word_wrap = True
tf.paragraphs[0].text = "💻 Frontend — TabSummarization"
tf.paragraphs[0].font.size = Pt(15)
tf.paragraphs[0].font.color.rgb = GREEN
tf.paragraphs[0].font.bold = True
p = tf.add_paragraph()
p.text = "React 19 + Ant Design 6 + Tailwind 4"
p.font.size = Pt(12)
p.font.color.rgb = GRAY
p.space_before = Pt(6)
for line in ["• Cột 1: Chọn doc CUAD / upload .txt", "• Cột 2: Chọn thuật toán + bật/tắt BART", "• Cột 3: Thanh tiến trình 4 bước"]:
    p = tf.add_paragraph()
    p.text = line
    p.font.size = Pt(11)
    p.font.color.rgb = GRAY
    p.space_before = Pt(4)

card = add_card(slide, Inches(6.8), Inches(4.4), Inches(5.7), Inches(2.4), PURPLE)
tf = card.text_frame
tf.word_wrap = True
tf.paragraphs[0].text = "📊 Panel kết quả"
tf.paragraphs[0].font.size = Pt(15)
tf.paragraphs[0].font.color.rgb = PURPLE
tf.paragraphs[0].font.bold = True
for line in ["• 4 thẻ thống kê: từ gốc, câu split, câu giữ, tỉ lệ nén",
             "• Danh sách câu trích kèm score",
             "• Bản viết lại BART (chi tiết từng chunk)",
             "• ROUGE-1/2/L + precision/recall",
             "• Nút xuất .txt"]:
    p = tf.add_paragraph()
    p.text = line
    p.font.size = Pt(11)
    p.font.color.rgb = GRAY
    p.space_before = Pt(4)


# ════════════════════════════════════════════════════════════════
# SLIDE 10 — Kiến trúc thư mục
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
set_bg(slide)
add_slide_number(slide, 10)
add_tag(slide, Inches(0.8), Inches(0.5), "CẤU TRÚC", ACCENT_L)

add_textbox(slide, Inches(0.8), Inches(1.0), Inches(11), Inches(0.6),
            "Kiến trúc thư mục dự án", font_size=28, color=WHITE, bold=True)

# Left tree
card = add_card(slide, Inches(0.8), Inches(1.8), Inches(5.5), Inches(5.0), ACCENT)
tf = card.text_frame
tf.word_wrap = True
tf.paragraphs[0].text = "text_sumarisation/"
tf.paragraphs[0].font.size = Pt(11)
tf.paragraphs[0].font.color.rgb = ACCENT_L
tf.paragraphs[0].font.bold = True
tf.paragraphs[0].font.name = 'Consolas'
tree_left = """├── data/
│   ├── full_contract_txt/    # 510 .txt
│   └── master_clauses.csv    # 83 cột → reference
│
├── backend/
│   ├── config/settings.py
│   ├── preprocess/        # ① loader+cleaner+splitter
│   ├── extractive/        # ② tfidf/textrank/kmeans
│   ├── abstractive/       # ③ chunker + BART
│   ├── evaluate/          # ④ ROUGE + BERTScore
│   ├── hybrid/pipeline.py # orchestrator
│   ├── training/          # fine-tune BART
│   └── app/               # Flask API"""
p = tf.add_paragraph()
p.text = tree_left
p.font.size = Pt(10)
p.font.color.rgb = GRAY
p.font.name = 'Consolas'
p.space_before = Pt(6)

# Right tree
card = add_card(slide, Inches(6.8), Inches(1.8), Inches(5.7), Inches(5.0), GREEN)
tf = card.text_frame
tf.word_wrap = True
tf.paragraphs[0].text = "├── frontend/"
tf.paragraphs[0].font.size = Pt(11)
tf.paragraphs[0].font.color.rgb = GREEN
tf.paragraphs[0].font.bold = True
tf.paragraphs[0].font.name = 'Consolas'
tree_right = """│   └── src/components/
│       TabSummarization.jsx
│
├── deploy/
│   ├── fabfile.py            # Fabric deploy
│   └── requirements*.txt
│
└── outputs/
    ├── eval/                 # ROUGE results
    ├── models/               # fine-tuned BART
    └── summary_cache/        # cached summaries

Frontend → axios /api/*
  → Vite proxy
  → Flask localhost:9020"""
p = tf.add_paragraph()
p.text = tree_right
p.font.size = Pt(10)
p.font.color.rgb = GRAY
p.font.name = 'Consolas'
p.space_before = Pt(6)


# ════════════════════════════════════════════════════════════════
# SLIDE 11 — Tham số cấu hình
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
set_bg(slide)
add_slide_number(slide, 11)
add_tag(slide, Inches(0.8), Inches(0.5), "CẤU HÌNH", ORANGE)

add_textbox(slide, Inches(0.8), Inches(1.0), Inches(11), Inches(0.6),
            "Phụ lục — Các tham số cấu hình chính", font_size=28, color=WHITE, bold=True)

# Table
card = add_card(slide, Inches(0.8), Inches(1.8), Inches(11.7), Inches(3.8), ACCENT)
tf = card.text_frame
tf.word_wrap = True

# Header
header = "  Tham số                Mặc định    Ý nghĩa                              Ảnh hưởng"
tf.paragraphs[0].text = header
tf.paragraphs[0].font.size = Pt(10)
tf.paragraphs[0].font.color.rgb = MUTED
tf.paragraphs[0].font.bold = True
tf.paragraphs[0].font.name = 'Consolas'

rows = [
    "  TOP_K_RATIO             0.2         Tỉ lệ câu extractive giữ lại         Cao → nhiều thông tin; Thấp → mất ý",
    "  MIN_SENT_LEN            5           Bỏ câu < N từ                        Loại tiêu đề ngắn, dòng nhiễu",
    "  MAX_SENT_LEN            80          Bỏ câu > N từ                        Loại đoạn dài bị tách nhầm",
    "  BART_MAX_INPUT          1024        Token tối đa BART nhận               Giới hạn cứng của mô hình",
    "  BART_MAX_OUTPUT         256         Token sinh ra tối đa                 → ~150–300 từ tóm tắt",
    "  BART_NUM_BEAMS          4           Beam search width                    Cao → chất lượng tốt hơn, chậm hơn",
    "  CHUNK_OVERLAP           50          Token chồng lấn giữa chunk           Giữ ngữ cảnh xuyên chunk",
]
for row in rows:
    p = tf.add_paragraph()
    p.text = row
    p.font.size = Pt(9)
    p.font.color.rgb = GRAY
    p.font.name = 'Consolas'
    p.space_before = Pt(3)

# Stat cards
add_stat_card(slide, Inches(1.2), Inches(5.9), "20%", "", "Câu giữ lại (extractive)", ACCENT_L)
add_stat_card(slide, Inches(4.8), Inches(5.9), "1024", "", "Token max input BART", PURPLE)
add_stat_card(slide, Inches(8.4), Inches(5.9), "4", "", "Beam search width", GREEN)


# ════════════════════════════════════════════════════════════════
# SLIDE 12 — Hạn chế & Hướng tiếp
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
set_bg(slide)
add_slide_number(slide, 12)
add_tag(slide, Inches(0.8), Inches(0.5), "HẠN CHẾ & HƯỚNG TIẾP", RED)

add_textbox(slide, Inches(0.8), Inches(1.0), Inches(11), Inches(0.6),
            "Hạn chế hiện tại & Hướng phát triển", font_size=28, color=WHITE, bold=True)

# Left - Hạn chế
card = add_card(slide, Inches(0.8), Inches(1.8), Inches(5.5), Inches(4.5), RED)
tf = card.text_frame
tf.word_wrap = True
tf.paragraphs[0].text = "⚠️ Hạn chế"
tf.paragraphs[0].font.size = Pt(16)
tf.paragraphs[0].font.color.rgb = RED
tf.paragraphs[0].font.bold = True
limitations = [
    "BART pretrained trên tin tức → phong cách hơi gượng với văn bản pháp lý",
    "Extractive không hiểu ngữ nghĩa sâu → có thể bỏ lỡ điều khoản quan trọng nếu từ khoá hiếm",
    "Chunking mất liên kết → overlap 50 token chưa đủ cho hợp đồng phức tạp",
    "Reference từ clause → không phải gold summary thực sự, chỉ là proxy",
]
for line in limitations:
    p = tf.add_paragraph()
    p.text = "• " + line
    p.font.size = Pt(12)
    p.font.color.rgb = GRAY
    p.space_before = Pt(8)

# Right - Hướng tiếp
card = add_card(slide, Inches(6.8), Inches(1.8), Inches(5.7), Inches(4.5), GREEN)
tf = card.text_frame
tf.word_wrap = True
tf.paragraphs[0].text = "🚀 Hướng tiếp"
tf.paragraphs[0].font.size = Pt(16)
tf.paragraphs[0].font.color.rgb = GREEN
tf.paragraphs[0].font.bold = True
directions = [
    "Thử LED / Longformer — mô hình hỗ trợ 16K+ token, xử lý trực tiếp văn bản dài",
    "Pegasus-BillSum — pretrained trên văn bản pháp luật",
    "Chấm điểm hậu kiểm bằng QA — dùng Question Answering để kiểm tra tóm tắt",
    "Legal-specific pretraining — tiếp tục fine-tune trên corpus pháp lý lớn hơn",
    "Multimodal — xử lý cả bảng biểu, sơ đồ trong hợp đồng",
]
for line in directions:
    p = tf.add_paragraph()
    p.text = "• " + line
    p.font.size = Pt(12)
    p.font.color.rgb = GRAY
    p.space_before = Pt(8)

# Thank you
card = add_card(slide, Inches(3.5), Inches(6.5), Inches(6.3), Inches(0.7), ACCENT)
tf = card.text_frame
tf.word_wrap = True
tf.paragraphs[0].text = "Cảm ơn! — Hệ thống Hybrid Summarization trên CUAD"
tf.paragraphs[0].font.size = Pt(14)
tf.paragraphs[0].font.color.rgb = WHITE
tf.paragraphs[0].font.bold = True
tf.paragraphs[0].alignment = PP_ALIGN.CENTER
tf.vertical_anchor = MSO_ANCHOR.MIDDLE


# ── Save ──
output_path = "/home/sontn/Projects/legal-ai/slides/presentation.pptx"
prs.save(output_path)
print(f"✅ Saved: {output_path}")